import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def convert_pdf_to_ppt(pdf_path, output_ppt, dpi=200):
    prs = Presentation()

    # Blank slide layout
    blank_layout = prs.slide_layouts[6]

    doc = fitz.open(pdf_path)
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    for page in doc:
        slide = prs.slides.add_slide(blank_layout)

        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = os.path.join(
            os.path.dirname(output_ppt),
            f"page_{page.number + 1}.png"
        )
        pix.save(img_path)

        slide.shapes.add_picture(
            img_path,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(output_ppt)
    doc.close()
